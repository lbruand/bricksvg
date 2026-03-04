"""svg2pptx — Convert an SVG file to an Office 2007 PPTX, one element per shape.

Each SVG element type becomes a native PPTX object:
  - <use> brick elements  → coloured PNG (grayscale render + duotone filter applied in PIL)
  - <line> grid elements  → connector shapes
  - <polygon> arrows      → rasterised PNG on transparent canvas
  - <text> labels         → DrawingML custom-geometry paths (cubic beziers from font outlines,
                            affine-transformed to preserve isometric perspective); falls back
                            to rasterised bitmap if no system font is found
  - <image> icons         → affine-warped PNG (preserves isometric top-face mapping)
"""

import argparse
import base64
import functools
import io
import math
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from lxml import etree as _et  # type: ignore[import-untyped]
from PIL import Image, ImageChops, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE


SVG_NS   = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"

_PX_TO_EMU = 9525  # 914 400 EMU/inch ÷ 96 px/inch

# DrawingML / PresentationML namespace strings used when building <p:sp> elements
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# System font search order (first existing file wins)
_FONT_CANDIDATES = [
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "C:/Windows/Fonts/arial.ttf",
]


# ---------------------------------------------------------------------------
# Unit helpers
# ---------------------------------------------------------------------------

def _emu(px: float) -> int:
    return round(px * _PX_TO_EMU)


def _parse_length(val: str, full: float = 0.0) -> float:
    """Parse an SVG length value; handles 'px' suffix and '%' relative lengths."""
    s = str(val).strip()
    if s.endswith("%"):
        return full * float(s[:-1]) / 100.0
    return float(s.replace("px", "").strip())


# ---------------------------------------------------------------------------
# SVG <defs> extraction
# ---------------------------------------------------------------------------

def _extract_defs(
    root: ET.Element,
) -> tuple[dict[str, Image.Image], dict[str, tuple[float, float, float]]]:
    """Return (images, filters) from SVG <defs>.

    images:  id → PIL RGBA image (the grayscale render stored in the SVG)
    filters: id → (r, g, b) floats 0–1 from the feColorMatrix duotone filter
    """
    images:  dict[str, Image.Image]              = {}
    filters: dict[str, tuple[float, float, float]] = {}

    defs = root.find(f"{{{SVG_NS}}}defs")
    if defs is None:
        return images, filters

    for img in defs.findall(f"{{{SVG_NS}}}image"):
        img_id = img.get("id", "")
        href   = img.get(f"{{{XLINK_NS}}}href") or img.get("href", "")
        if href.startswith("data:image/png;base64,"):
            raw = base64.b64decode(href[len("data:image/png;base64,"):])
            images[img_id] = Image.open(io.BytesIO(raw)).convert("RGBA")

    for filt in defs:
        if filt.tag != f"{{{SVG_NS}}}filter":
            continue
        fid = filt.get("id", "")
        for child in filt:
            if child.tag == f"{{{SVG_NS}}}feColorMatrix":
                nums = [float(x) for x in child.get("values", "").split()]
                # Our matrix: r 0 0 0 0  g 0 0 0 0  b 0 0 0 0  0 0 0 1 0
                # → nums[0]=r, nums[5]=g, nums[10]=b
                if len(nums) >= 11:
                    filters[fid] = (nums[0], nums[5], nums[10])

    return images, filters


# ---------------------------------------------------------------------------
# PIL helpers
# ---------------------------------------------------------------------------

def _apply_duotone(
    img: Image.Image,
    rgb_f: tuple[float, float, float],
) -> Image.Image:
    """Replicate feColorMatrix multiply: white → brick colour, shadows darken."""
    r, g, b = (round(c * 255) for c in rgb_f)
    solid   = Image.new("RGB", img.size, (r, g, b))
    result  = ImageChops.multiply(img.convert("RGB"), solid)
    result.putalpha(img.split()[3])
    return result.convert("RGBA")


def _to_stream(img: Image.Image) -> io.BytesIO:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _add_picture(slide, img: Image.Image, x: float, y: float) -> None:
    iw, ih = img.size
    slide.shapes.add_picture(_to_stream(img), _emu(x), _emu(y), _emu(iw), _emu(ih))


# ---------------------------------------------------------------------------
# Per-element processors
# ---------------------------------------------------------------------------

def _process_use(
    el: ET.Element,
    images: dict[str, Image.Image],
    filters: dict[str, tuple[float, float, float]],
    cache: dict[tuple[str, str], Image.Image],
    slide,
) -> None:
    href    = el.get(f"{{{XLINK_NS}}}href") or el.get("href", "")
    img_id  = href.lstrip("#")
    filt_id = re.sub(r"url\(#(.+)\)", r"\1", el.get("filter", ""))
    x       = _parse_length(el.get("x", "0"))
    y       = _parse_length(el.get("y", "0"))

    base = images.get(img_id)
    if base is None:
        return

    key = (img_id, filt_id)
    if key not in cache:
        rgb_f      = filters.get(filt_id)
        cache[key] = _apply_duotone(base, rgb_f) if rgb_f else base
    _add_picture(slide, cache[key], x, y)


def _process_line(el: ET.Element, slide) -> None:
    x1 = _parse_length(el.get("x1", "0"))
    y1 = _parse_length(el.get("y1", "0"))
    x2 = _parse_length(el.get("x2", "0"))
    y2 = _parse_length(el.get("y2", "0"))

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        _emu(x1), _emu(y1), _emu(x2), _emu(y2),
    )
    stroke = el.get("stroke", "#cccccc")
    if stroke.startswith("#") and len(stroke) == 7:
        conn.line.color.rgb = RGBColor(
            int(stroke[1:3], 16), int(stroke[3:5], 16), int(stroke[5:7], 16),
        )
    try:
        conn.line.width = _emu(_parse_length(el.get("stroke-width", "1")))
    except Exception:
        pass


def _parse_hex_color(val: str) -> tuple[int, int, int]:
    """Parse a 3- or 6-digit SVG hex colour string into (r, g, b)."""
    h = val.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _process_polygon(
    el: ET.Element,
    slide,
    svg_w: float,
    svg_h: float,
) -> None:
    pts_str = el.get("points", "")
    nums    = [float(v) for v in re.split(r"[,\s]+", pts_str.strip()) if v]
    pts     = [(nums[i], nums[i + 1]) for i in range(0, len(nums) - 1, 2)]
    if not pts:
        return

    r, g, b = _parse_hex_color(el.get("fill", "#888888"))
    color   = (r, g, b, 220)

    xs, ys = [p[0] for p in pts], [p[1] for p in pts]
    pad = 2.0
    x0  = max(0.0, min(xs) - pad)
    y0  = max(0.0, min(ys) - pad)
    x1  = min(svg_w, max(xs) + pad)
    y1  = min(svg_h, max(ys) + pad)

    canvas = Image.new("RGBA", (round(x1 - x0), round(y1 - y0)), (0, 0, 0, 0))
    ImageDraw.Draw(canvas).polygon([(px - x0, py - y0) for px, py in pts], fill=color)
    _add_picture(slide, canvas, x0, y0)


def _parse_matrix(transform: str) -> tuple[float, float, float, float, float, float]:
    """Return (a,b,c,d,e,f) from a SVG matrix(...) transform, or the identity."""
    m = re.search(r"matrix\(([^)]+)\)", transform)
    if not m:
        return 1.0, 0.0, 0.0, 1.0, 0.0, 0.0
    vals = [float(v) for v in re.split(r"[,\s]+", m.group(1).strip()) if v]
    return (vals[0], vals[1], vals[2], vals[3], vals[4], vals[5]) if len(vals) >= 6 \
        else (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)


def _parse_matrix_translation(transform: str) -> tuple[float, float]:
    """Extract (tx, ty) from a SVG matrix(a,b,c,d,e,f) transform."""
    _, _, _, _, e, f = _parse_matrix(transform)
    return e, f


def _affine_warp(
    img: Image.Image,
    a: float, b: float, c: float, d: float, e: float, f: float,
) -> tuple[Image.Image, float, float]:
    """Warp *img* with a SVG forward matrix; return (result, canvas_x, canvas_y).

    SVG forward mapping:  screen_x = a*src_x + c*src_y + e
                          screen_y = b*src_x + d*src_y + f

    Returns the warped image and the screen-space top-left corner at which to
    place it so the visual result matches the original SVG element.
    """
    W, H = img.size
    corners = [(a*sx + c*sy + e, b*sx + d*sy + f)
               for sx, sy in [(0, 0), (W, 0), (0, H), (W, H)]]
    xs = [p[0] for p in corners]
    ys = [p[1] for p in corners]
    cx0, cy0 = min(xs), min(ys)
    cw = math.ceil(max(xs) - cx0)
    ch = math.ceil(max(ys) - cy0)
    det = a * d - b * c
    if abs(det) < 1e-10 or cw <= 0 or ch <= 0:
        return img, e, f

    # PIL AFFINE (inverse): canvas pixel (px, py) → source pixel (src_x, src_y).
    # Canvas origin is at SVG position (cx0, cy0), so SVG = canvas + (cx0, cy0).
    # Solving: [SVG - translation] = M2 * src  →  src = M2_inv * (canvas + offset)
    ex, fy = cx0 - e, cy0 - f
    pil_data = (
        d / det,  -c / det, (d * ex - c * fy) / det,
        -b / det,  a / det, (-b * ex + a * fy) / det,
    )
    warped = img.transform(
        (cw, ch), Image.Transform.AFFINE, pil_data,
        resample=Image.Resampling.BICUBIC,
    )
    return warped, cx0, cy0


def _rasterize_text(
    text: str,
    font_size: float,
    fill_rgb: tuple[int, int, int],
) -> Image.Image:
    """Render *text* to a transparent RGBA image (fallback when no vector font found)."""
    font  = ImageFont.load_default(size=round(max(8, font_size)))
    probe = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    bbox  = probe.textbbox((0, 0), text or " ", font=font)
    tw    = max(1, int(bbox[2] - bbox[0]))
    th    = max(1, int(bbox[3] - bbox[1]))
    img   = Image.new("RGBA", (tw, th), (0, 0, 0, 0))
    ImageDraw.Draw(img).text((-bbox[0], -bbox[1]), text or " ", font=font,
                             fill=fill_rgb + (255,))
    return img


# ---------------------------------------------------------------------------
# Vector text: font outlines → DrawingML custom-geometry paths
# ---------------------------------------------------------------------------

@functools.lru_cache(maxsize=1)
def _load_font():
    """Return a fonttools TTFont for the first available system font, or None."""
    from fontTools.ttLib import TTFont
    for path in _FONT_CANDIDATES:
        if Path(path).exists():
            try:
                return TTFont(path)
            except Exception:
                continue
    return None


def _quad_to_cubic(
    p0: tuple[float, float],
    p1: tuple[float, float],
    p2: tuple[float, float],
) -> tuple[tuple[float, float], tuple[float, float], tuple[float, float]]:
    """Convert a quadratic bezier (p0→p1→p2) to cubic control points."""
    return (
        (p0[0] + 2/3*(p1[0]-p0[0]), p0[1] + 2/3*(p1[1]-p0[1])),
        (p2[0] + 2/3*(p1[0]-p2[0]), p2[1] + 2/3*(p1[1]-p2[1])),
        p2,
    )


def _text_to_ops_emu(
    text: str,
    font_size: float,
    a: float, b: float, c: float, d: float, e: float, f: float,
) -> list | None:
    """Convert *text* to DrawingML path operations in EMU with SVG affine applied.

    Returns a list of ``(op, pts)`` tuples where *op* is one of
    ``'moveTo'``, ``'lineTo'``, ``'curveTo'`` (cubic), ``'closePath'``;
    *pts* is a list of ``(x_emu, y_emu)`` integers.
    Returns ``None`` if no system font is available.

    Glyphs are horizontally centred on the transform origin
    (matching SVG ``text_anchor="middle"``).
    """
    from fontTools.pens.recordingPen import RecordingPen

    font = _load_font()
    if font is None:
        return None

    upm       = font["head"].unitsPerEm
    scale     = font_size / upm          # font design units → pixels
    glyph_set = font.getGlyphSet()
    cmap      = font.getBestCmap() or {}

    # Pass 1: resolve glyph names and accumulate advance widths for centering
    char_glyphs: list[tuple] = []
    x_cursor = 0.0
    for ch in text:
        gname = cmap.get(ord(ch))
        if gname and gname in glyph_set:
            g = glyph_set[gname]
            char_glyphs.append((g, x_cursor))
            x_cursor += g.width
        else:
            x_cursor += upm * 0.55          # fallback advance for missing glyphs

    x_origin = -x_cursor / 2               # centre on the transform anchor

    def to_px(fx: float, fy: float, x_off: float) -> tuple[float, float]:
        """Font units (Y-up) → pixels (Y-down), shifted by glyph x_off and centring."""
        return (fx + x_off + x_origin) * scale, -fy * scale

    def transform_emu(px: float, py: float) -> tuple[int, int]:
        """Apply SVG matrix and convert to EMU."""
        return (round((a*px + c*py + e) * _PX_TO_EMU),
                round((b*px + d*py + f) * _PX_TO_EMU))

    all_ops: list = []
    for g, x_off in char_glyphs:
        pen = RecordingPen()
        g.draw(pen)

        current: tuple[float, float] = (0.0, 0.0)   # last on-curve point in pixels

        for op_name, args in pen.value:
            if op_name == "moveTo":
                p = to_px(args[0][0], args[0][1], x_off)
                all_ops.append(("moveTo", [transform_emu(*p)]))
                current = p

            elif op_name == "lineTo":
                p = to_px(args[0][0], args[0][1], x_off)
                all_ops.append(("lineTo", [transform_emu(*p)]))
                current = p

            elif op_name == "curveTo":
                # Cubic bezier: 2 control points + endpoint
                pts = [to_px(pt[0], pt[1], x_off) for pt in args]
                all_ops.append(("curveTo", [transform_emu(*pt) for pt in pts]))
                current = pts[-1]

            elif op_name == "qCurveTo":
                # TrueType quadratic: all but last are off-curve, last is on-curve.
                # Multiple consecutive off-curves imply on-curves at their midpoints.
                offs = [to_px(pt[0], pt[1], x_off) for pt in args[:-1]]
                end  = to_px(args[-1][0], args[-1][1], x_off)
                prev = current
                for i, off in enumerate(offs):
                    nxt = (
                        ((offs[i][0] + offs[i+1][0]) / 2,
                         (offs[i][1] + offs[i+1][1]) / 2)
                        if i + 1 < len(offs) else end
                    )
                    c1, c2, ep = _quad_to_cubic(prev, off, nxt)
                    all_ops.append(("curveTo", [transform_emu(*pt) for pt in (c1, c2, ep)]))
                    prev = nxt
                current = end

            elif op_name in ("closePath", "endPath"):
                all_ops.append(("closePath", []))

    return all_ops


def _build_text_sp(
    ops_emu: list,
    fill_rgb: tuple[int, int, int],
) -> "_et._Element | None":
    """Build a DrawingML ``<p:sp>`` with custom-geometry paths from *ops_emu*.

    All path coordinates in *ops_emu* are absolute EMU values.  The shape
    bounding box is computed from the point set; path coords inside the element
    are relative to the shape's top-left corner.
    """
    all_pts = [pt for op, pts in ops_emu for pt in pts]
    if not all_pts:
        return None

    x0 = min(p[0] for p in all_pts)
    y0 = min(p[1] for p in all_pts)
    x1 = max(p[0] for p in all_pts)
    y1 = max(p[1] for p in all_pts)
    cx = max(1, x1 - x0)
    cy = max(1, y1 - y0)

    def _a(tag: str) -> str: return f"{{{_A_NS}}}{tag}"
    def _p(tag: str) -> str: return f"{{{_P_NS}}}{tag}"
    def _pt(parent, x: int, y: int) -> None:
        e = _et.SubElement(parent, _a("pt"))
        e.set("x", str(x - x0))
        e.set("y", str(y - y0))

    sp = _et.Element(_p("sp"))

    # Non-visual properties
    nvSpPr  = _et.SubElement(sp, _p("nvSpPr"))
    cNvPr   = _et.SubElement(nvSpPr, _p("cNvPr"))
    cNvPr.set("id", "1")
    cNvPr.set("name", "text")
    cNvSpPr = _et.SubElement(nvSpPr, _p("cNvSpPr"))
    _et.SubElement(cNvSpPr, _a("spLocks")).set("noGrp", "1")
    _et.SubElement(nvSpPr, _p("nvPr"))

    # Shape properties
    spPr = _et.SubElement(sp, _p("spPr"))
    xfrm = _et.SubElement(spPr, _a("xfrm"))
    off  = _et.SubElement(xfrm, _a("off"))
    off.set("x", str(x0))
    off.set("y", str(y0))
    ext  = _et.SubElement(xfrm, _a("ext"))
    ext.set("cx", str(cx))
    ext.set("cy", str(cy))

    # Custom geometry
    custGeom = _et.SubElement(spPr, _a("custGeom"))
    for tag in ("avLst", "gdLst", "ahLst", "cxnLst"):
        _et.SubElement(custGeom, _a(tag))
    rect = _et.SubElement(custGeom, _a("rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(cx))
    rect.set("b", str(cy))

    pathLst = _et.SubElement(custGeom, _a("pathLst"))
    path_el = _et.SubElement(pathLst, _a("path"))
    path_el.set("w", str(cx))
    path_el.set("h", str(cy))

    for op, pts in ops_emu:
        if op == "moveTo":
            _pt(_et.SubElement(path_el, _a("moveTo")), *pts[0])
        elif op == "lineTo":
            _pt(_et.SubElement(path_el, _a("lnTo")), *pts[0])
        elif op == "curveTo":
            cb = _et.SubElement(path_el, _a("cubicBezTo"))
            for pt in pts:
                _pt(cb, *pt)
        elif op == "closePath":
            _et.SubElement(path_el, _a("close"))

    # Solid fill
    srgb = _et.SubElement(_et.SubElement(spPr, _a("solidFill")), _a("srgbClr"))
    srgb.set("val", "{:02x}{:02x}{:02x}".format(*fill_rgb))

    # No outline stroke
    _et.SubElement(_et.SubElement(spPr, _a("ln")), _a("noFill"))

    return sp


def _add_sp(slide, sp_el: "_et._Element") -> None:
    """Append *sp_el* to the slide shape tree, assigning a unique shape ID."""
    spTree = slide.shapes._spTree
    cNvPr_tag = f"{{{_P_NS}}}cNvPr"
    max_id = max(
        (int(el.get("id", 0)) for el in spTree.iter(cNvPr_tag)),
        default=0,
    )
    cNvPr = sp_el.find(f".//{cNvPr_tag}")
    if cNvPr is not None:
        cNvPr.set("id", str(max_id + 1))
    spTree.append(sp_el)


def _process_text(el: ET.Element, slide) -> None:
    a, b, c, d, e, f = _parse_matrix(el.get("transform", ""))
    text = el.text or ""
    fill = el.get("fill", "#000000")
    try:
        fs = float(el.get("font-size") or el.get("font_size") or "14")
    except ValueError:
        fs = 14.0
    fill_rgb = _parse_hex_color(fill) if fill.startswith("#") else (0, 0, 0)

    ops = _text_to_ops_emu(text, fs, a, b, c, d, e, f)
    if ops is not None:
        sp = _build_text_sp(ops, fill_rgb)
        if sp is not None:
            _add_sp(slide, sp)
            return

    # Fallback: rasterise + affine-warp bitmap
    img = _rasterize_text(text, fs, fill_rgb)
    W, H = img.size
    warped, wx, wy = _affine_warp(img, a, b, c, d,
                                   e - a*W/2 - c*H/2,
                                   f - b*W/2 - d*H/2)
    _add_picture(slide, warped, wx, wy)


def _process_icon(el: ET.Element, slide) -> None:
    href = el.get(f"{{{XLINK_NS}}}href") or el.get("href", "")
    if not href.startswith("data:image/png;base64,"):
        return
    raw = base64.b64decode(href[len("data:image/png;base64,"):])
    img = Image.open(io.BytesIO(raw)).convert("RGBA")
    a, b, c, d, e, f = _parse_matrix(el.get("transform", ""))
    warped, wx, wy = _affine_warp(img, a, b, c, d, e, f)
    _add_picture(slide, warped, wx, wy)


# ---------------------------------------------------------------------------
# SVG tree walker
# ---------------------------------------------------------------------------

def _walk(
    root: ET.Element,
    images: dict[str, Image.Image],
    filters: dict[str, tuple[float, float, float]],
    cache: dict[tuple[str, str], Image.Image],
    slide,
    svg_w: float,
    svg_h: float,
    *,
    depth: int = 0,
) -> None:
    """Recursively process SVG elements, placing each as a slide object."""
    for el in root:
        tag = el.tag
        if tag == f"{{{SVG_NS}}}defs":
            continue
        elif tag == f"{{{SVG_NS}}}rect" and depth == 0:
            pass  # background handled via slide.background.fill
        elif tag == f"{{{SVG_NS}}}use":
            _process_use(el, images, filters, cache, slide)
        elif tag == f"{{{SVG_NS}}}line":
            _process_line(el, slide)
        elif tag == f"{{{SVG_NS}}}polygon":
            _process_polygon(el, slide, svg_w, svg_h)
        elif tag == f"{{{SVG_NS}}}text":
            _process_text(el, slide)
        elif tag == f"{{{SVG_NS}}}image":
            _process_icon(el, slide)
        elif tag == f"{{{SVG_NS}}}g":
            _walk(el, images, filters, cache, slide, svg_w, svg_h, depth=depth + 1)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert an SVG file to an Office 2007 PPTX, one element per shape.",
    )
    parser.add_argument("input",  help="Input SVG file")
    parser.add_argument("-o", "--output", help="Output PPTX file (default: <input>.pptx)")
    args = parser.parse_args()

    input_path  = Path(args.input)
    output_path = Path(args.output or input_path.with_suffix(".pptx"))

    print(f"Parsing {input_path} …")
    root  = ET.parse(input_path).getroot()
    svg_w = _parse_length(root.get("width",  "800px"))
    svg_h = _parse_length(root.get("height", "600px"))
    print(f"  SVG: {svg_w:.0f} × {svg_h:.0f} px")

    images, filters = _extract_defs(root)
    print(f"  Defs: {len(images)} images, {len(filters)} filters")

    prs              = Presentation()
    prs.slide_width  = _emu(svg_w)   # type: ignore[assignment]
    prs.slide_height = _emu(svg_h)   # type: ignore[assignment]
    slide            = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Slide background colour from the top-level <rect>
    bg_rect = root.find(f"{{{SVG_NS}}}rect")
    if bg_rect is not None:
        fill = bg_rect.get("fill", "#f8f8f0")
        if fill.startswith("#") and len(fill) == 7:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(
                int(fill[1:3], 16), int(fill[3:5], 16), int(fill[5:7], 16),
            )

    cache: dict[tuple[str, str], Image.Image] = {}
    print("Placing elements …")
    _walk(root, images, filters, cache, slide, svg_w, svg_h)

    n_shapes = len(slide.shapes)
    prs.save(str(output_path))
    print(f"Saved: {output_path}  ({n_shapes} shapes, {len(cache)} cached renders)")


if __name__ == "__main__":
    main()
