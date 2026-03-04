"""Unit tests for ldr2svg.svg2pptx."""

import base64
import io
import sys
import xml.etree.ElementTree as ET
from pathlib import Path
from unittest.mock import patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor

from ldr2svg.svg2pptx import (
    SVG_NS,
    XLINK_NS,
    _affine_warp,
    _apply_duotone,
    _emu,
    _extract_defs,
    _parse_hex_color,
    _parse_length,
    _parse_matrix,
    _parse_matrix_translation,
    _process_icon,
    _process_line,
    _process_polygon,
    _process_text,
    _process_use,
    _rasterize_text,
    _to_stream,
    _walk,
    main,
)

DIAGRAM_SVG = Path(__file__).parent.parent / "diagram_bricks.svg"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _png_data_uri(w=10, h=10, color=(128, 128, 128, 255)):
    img = Image.new("RGBA", (w, h), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode()
    return f"data:image/png;base64,{b64}"


# ---------------------------------------------------------------------------
# TestEmu
# ---------------------------------------------------------------------------

class TestEmu:
    def test_zero(self):
        assert _emu(0) == 0

    def test_one_pixel(self):
        assert _emu(1) == 9525

    def test_hundred_pixels(self):
        assert _emu(100) == 952500

    def test_returns_int(self):
        assert isinstance(_emu(1.7), int)


# ---------------------------------------------------------------------------
# TestParseLength
# ---------------------------------------------------------------------------

class TestParseLength:
    def test_bare_number(self):
        assert _parse_length("100") == pytest.approx(100.0)

    def test_px_suffix(self):
        assert _parse_length("200px") == pytest.approx(200.0)

    def test_float_string(self):
        assert _parse_length("12.5") == pytest.approx(12.5)

    def test_percent_of_full(self):
        assert _parse_length("50%", full=800.0) == pytest.approx(400.0)

    def test_percent_zero_full(self):
        assert _parse_length("100%", full=0.0) == pytest.approx(0.0)


# ---------------------------------------------------------------------------
# TestParseHexColor
# ---------------------------------------------------------------------------

class TestParseHexColor:
    def test_six_digit(self):
        assert _parse_hex_color("#ff8800") == (255, 136, 0)

    def test_black(self):
        assert _parse_hex_color("#000000") == (0, 0, 0)

    def test_white(self):
        assert _parse_hex_color("#ffffff") == (255, 255, 255)

    def test_three_digit_expands_each_nibble(self):
        assert _parse_hex_color("#abc") == (0xaa, 0xbb, 0xcc)

    def test_three_digit_gray(self):
        assert _parse_hex_color("#888") == (136, 136, 136)


# ---------------------------------------------------------------------------
# TestParseMatrixTranslation
# ---------------------------------------------------------------------------

class TestParseMatrixTranslation:
    def test_comma_separated(self):
        tx, ty = _parse_matrix_translation("matrix(1,0,0,1,42,99)")
        assert tx == pytest.approx(42.0)
        assert ty == pytest.approx(99.0)

    def test_space_separated(self):
        tx, ty = _parse_matrix_translation("matrix(1 0 0 1 10 20)")
        assert tx == pytest.approx(10.0)
        assert ty == pytest.approx(20.0)

    def test_float_values(self):
        tx, ty = _parse_matrix_translation("matrix(0.866,0.5,-0.866,0.5,123.4,56.7)")
        assert tx == pytest.approx(123.4)
        assert ty == pytest.approx(56.7)

    def test_empty_string(self):
        assert _parse_matrix_translation("") == (0.0, 0.0)

    def test_non_matrix_transform(self):
        assert _parse_matrix_translation("translate(5,10)") == (0.0, 0.0)


# ---------------------------------------------------------------------------
# TestApplyDuotone
# ---------------------------------------------------------------------------

class TestParseMatrix:
    def test_comma_separated(self):
        a, b, c, d, e, f = _parse_matrix("matrix(0.866,0.5,-0.866,0.5,100,200)")
        assert (a, b, c, d, e, f) == pytest.approx((0.866, 0.5, -0.866, 0.5, 100, 200))

    def test_space_separated(self):
        a, b, c, d, e, f = _parse_matrix("matrix(1 0 0 1 42 99)")
        assert (e, f) == pytest.approx((42.0, 99.0))

    def test_identity_fallback_on_empty(self):
        assert _parse_matrix("") == (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)

    def test_identity_fallback_on_non_matrix(self):
        assert _parse_matrix("translate(10,20)") == (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)

    def test_translation_components(self):
        _, _, _, _, e, f = _parse_matrix("matrix(1,0,0,1,42,99)")
        assert e == pytest.approx(42.0)
        assert f == pytest.approx(99.0)


class TestAffineWarp:
    def test_identity_returns_same_size(self):
        img = Image.new("RGBA", (30, 20))
        warped, _, _ = _affine_warp(img, 1.0, 0.0, 0.0, 1.0, 0.0, 0.0)
        assert warped.size == (30, 20)

    def test_identity_positions_at_translation(self):
        img = Image.new("RGBA", (30, 20))
        _, x, y = _affine_warp(img, 1.0, 0.0, 0.0, 1.0, 50.0, 100.0)
        assert x == pytest.approx(50.0)
        assert y == pytest.approx(100.0)

    def test_isometric_widens_bounding_box(self):
        """Left-face matrix matrix(0.866,0.5,0,1,...) produces a taller canvas."""
        img = Image.new("RGBA", (40, 40))
        warped, _, _ = _affine_warp(img, 0.866025, 0.5, 0.0, 1.0, 0.0, 0.0)
        # Height = 0.5*W + H = 20 + 40 = 60; width = 0.866*W ≈ 35
        assert warped.size[1] > 40

    def test_top_face_widens_canvas(self):
        """Top-face matrix(0.866,0.5,-0.866,0.5,...) widens the bounding box."""
        img = Image.new("RGBA", (40, 40))
        warped, _, _ = _affine_warp(img, 0.866025, 0.5, -0.866025, 0.5, 0.0, 0.0)
        assert warped.size[0] > 40

    def test_degenerate_matrix_returns_original_position(self):
        img = Image.new("RGBA", (10, 10))
        _, x, y = _affine_warp(img, 0.0, 0.0, 0.0, 0.0, 5.0, 10.0)
        assert (x, y) == (5.0, 10.0)

    def test_output_is_rgba(self):
        img = Image.new("RGBA", (20, 20), (255, 0, 0, 255))
        warped, _, _ = _affine_warp(img, 0.866025, 0.5, 0.0, 1.0, 0.0, 0.0)
        assert warped.mode == "RGBA"


class TestRasterizeText:
    def test_returns_rgba_image(self):
        img = _rasterize_text("Hello", 14, (0, 0, 0))
        assert isinstance(img, Image.Image)
        assert img.mode == "RGBA"

    def test_non_zero_size(self):
        img = _rasterize_text("Hello", 14, (0, 0, 0))
        assert img.size[0] > 0 and img.size[1] > 0

    def test_larger_font_produces_larger_image(self):
        small = _rasterize_text("A", 12, (0, 0, 0))
        large = _rasterize_text("A", 36, (0, 0, 0))
        assert large.size[0] >= small.size[0]

    def test_empty_text_handled(self):
        img = _rasterize_text("", 14, (0, 0, 0))
        assert isinstance(img, Image.Image)


class TestApplyDuotone:
    def test_white_becomes_target_color(self):
        img = Image.new("RGBA", (4, 4), (255, 255, 255, 255))
        result = _apply_duotone(img, (1.0, 0.0, 0.0))
        r, g, b, _ = result.getpixel((0, 0))
        assert r == 255
        assert g == 0
        assert b == 0

    def test_black_stays_black(self):
        img = Image.new("RGBA", (4, 4), (0, 0, 0, 255))
        result = _apply_duotone(img, (1.0, 0.5, 0.0))
        r, g, b, _ = result.getpixel((0, 0))
        assert r == 0 and g == 0 and b == 0

    def test_alpha_preserved(self):
        img = Image.new("RGBA", (4, 4), (255, 255, 255, 128))
        result = _apply_duotone(img, (0.5, 0.5, 0.5))
        _, _, _, a = result.getpixel((0, 0))
        assert a == 128

    def test_output_mode_is_rgba(self):
        img = Image.new("RGBA", (4, 4), (255, 255, 255, 255))
        result = _apply_duotone(img, (0.5, 0.5, 0.5))
        assert result.mode == "RGBA"


# ---------------------------------------------------------------------------
# TestToStream
# ---------------------------------------------------------------------------

class TestToStream:
    def test_returns_bytesio(self):
        assert isinstance(_to_stream(Image.new("RGBA", (10, 10))), io.BytesIO)

    def test_seek_position_is_zero(self):
        stream = _to_stream(Image.new("RGBA", (5, 5)))
        assert stream.tell() == 0

    def test_decodable_back_to_image(self):
        img = Image.new("RGBA", (10, 10), (100, 150, 200, 255))
        decoded = Image.open(_to_stream(img))
        assert decoded.size == (10, 10)


# ---------------------------------------------------------------------------
# TestExtractDefs
# ---------------------------------------------------------------------------

class TestExtractDefs:
    def _svg_with_defs(self, img_id="img1", filt_id="f1"):
        svg = ET.Element(f"{{{SVG_NS}}}svg")
        defs = ET.SubElement(svg, f"{{{SVG_NS}}}defs")
        img_el = ET.SubElement(defs, f"{{{SVG_NS}}}image")
        img_el.set("id", img_id)
        img_el.set("href", _png_data_uri())
        filt_el = ET.SubElement(defs, f"{{{SVG_NS}}}filter")
        filt_el.set("id", filt_id)
        cm = ET.SubElement(filt_el, f"{{{SVG_NS}}}feColorMatrix")
        cm.set("values", "0.5 0 0 0 0  0.25 0 0 0 0  0.1 0 0 0 0  0 0 0 1 0")
        return svg

    def test_image_key_present(self):
        images, _ = _extract_defs(self._svg_with_defs())
        assert "img1" in images

    def test_image_is_pil_rgba(self):
        images, _ = _extract_defs(self._svg_with_defs())
        assert isinstance(images["img1"], Image.Image)
        assert images["img1"].mode == "RGBA"

    def test_filter_key_present(self):
        _, filters = _extract_defs(self._svg_with_defs())
        assert "f1" in filters

    def test_filter_values_correct(self):
        _, filters = _extract_defs(self._svg_with_defs())
        r, g, b = filters["f1"]
        assert r == pytest.approx(0.5)
        assert g == pytest.approx(0.25)
        assert b == pytest.approx(0.1)

    def test_no_defs_returns_empty(self):
        svg = ET.Element(f"{{{SVG_NS}}}svg")
        images, filters = _extract_defs(svg)
        assert images == {} and filters == {}


# ---------------------------------------------------------------------------
# TestProcessUse
# ---------------------------------------------------------------------------

class TestProcessUse:
    def _make_el(self, img_id="img1", filt_id="f1", x=10.0, y=20.0):
        el = ET.Element(f"{{{SVG_NS}}}use")
        el.set("href", f"#{img_id}")
        el.set("filter", f"url(#{filt_id})")
        el.set("x", str(x))
        el.set("y", str(y))
        return el

    def _base_img(self):
        return Image.new("RGBA", (30, 20), (200, 200, 200, 255))

    def test_adds_one_shape(self):
        slide = _blank_slide()
        _process_use(self._make_el(), {"img1": self._base_img()},
                     {"f1": (1.0, 0.0, 0.0)}, {}, slide)
        assert len(slide.shapes) == 1

    def test_second_call_uses_cache(self):
        slide = _blank_slide()
        images = {"img1": self._base_img()}
        filters = {"f1": (0.5, 0.5, 0.5)}
        cache = {}
        el = self._make_el()
        _process_use(el, images, filters, cache, slide)
        _process_use(el, images, filters, cache, slide)
        assert len(cache) == 1

    def test_unknown_image_skipped(self):
        slide = _blank_slide()
        _process_use(self._make_el(), {}, {}, {}, slide)
        assert len(slide.shapes) == 0

    def test_no_filter_uses_base_image(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}use")
        el.set("href", "#img1")
        el.set("x", "0")
        el.set("y", "0")
        _process_use(el, {"img1": self._base_img()}, {}, {}, slide)
        assert len(slide.shapes) == 1

    def test_xlink_href_resolved(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}use")
        el.set(f"{{{XLINK_NS}}}href", "#img1")
        el.set("x", "0")
        el.set("y", "0")
        _process_use(el, {"img1": self._base_img()}, {}, {}, slide)
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# TestProcessLine
# ---------------------------------------------------------------------------

class TestProcessLine:
    def _make_el(self, x1=0, y1=0, x2=100, y2=50,
                 stroke="#ff0000", width="2"):
        el = ET.Element(f"{{{SVG_NS}}}line")
        el.set("x1", str(x1))
        el.set("y1", str(y1))
        el.set("x2", str(x2))
        el.set("y2", str(y2))
        el.set("stroke", stroke)
        el.set("stroke-width", width)
        return el

    def test_adds_connector(self):
        slide = _blank_slide()
        _process_line(self._make_el(), slide)
        assert len(slide.shapes) == 1

    def test_connector_color(self):
        slide = _blank_slide()
        _process_line(self._make_el(stroke="#00ff00"), slide)
        assert slide.shapes[0].line.color.rgb == RGBColor(0, 255, 0)

    def test_no_stroke_attribute_accepted(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}line")
        el.set("x1", "0")
        el.set("y1", "0")
        el.set("x2", "10")
        el.set("y2", "10")
        _process_line(el, slide)
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# TestProcessPolygon
# ---------------------------------------------------------------------------

class TestProcessPolygon:
    def test_adds_picture(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}polygon")
        el.set("points", "10,10 20,10 15,20")
        el.set("fill", "#ff8800")
        _process_polygon(el, slide, 800.0, 600.0)
        assert len(slide.shapes) == 1

    def test_three_digit_hex_fill(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}polygon")
        el.set("points", "10,10 20,10 15,20")
        el.set("fill", "#f80")
        _process_polygon(el, slide, 800.0, 600.0)
        assert len(slide.shapes) == 1

    def test_empty_points_skipped(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}polygon")
        el.set("points", "")
        _process_polygon(el, slide, 800.0, 600.0)
        assert len(slide.shapes) == 0


# ---------------------------------------------------------------------------
# TestProcessText
# ---------------------------------------------------------------------------

class TestProcessText:
    def _make_el(self, text="Hello", tx=100.0, ty=200.0,
                 fill="#333333", fs="14"):
        el = ET.Element(f"{{{SVG_NS}}}text")
        el.set("transform", f"matrix(1,0,0,1,{tx},{ty})")
        el.set("fill", fill)
        el.set("font-size", fs)
        el.text = text
        return el

    def test_adds_shape(self):
        slide = _blank_slide()
        _process_text(self._make_el(), slide)
        assert len(slide.shapes) == 1

    def test_isometric_left_face_adds_shape(self):
        """Text with isometric left-face matrix still adds a shape."""
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}text")
        el.set("transform", "matrix(0.866025,0.5,0,1,100,200)")
        el.set("fill", "#333333")
        el.set("font-size", "14")
        el.text = "AWS Lambda"
        _process_text(el, slide)
        assert len(slide.shapes) == 1

    def test_isometric_top_face_adds_shape(self):
        """Cluster labels use matrix(0.866025,0.5,-0.866025,0.5,...)."""
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}text")
        el.set("transform", "matrix(0.866025,0.5,-0.866025,0.5,200,150)")
        el.set("fill", "#444")
        el.set("font-size", "40")
        el.text = "my-cluster"
        _process_text(el, slide)
        assert len(slide.shapes) == 1

    def test_none_text_handled(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}text")
        el.set("transform", "matrix(1,0,0,1,10,20)")
        el.text = None
        _process_text(el, slide)
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# TestProcessIcon
# ---------------------------------------------------------------------------

class TestProcessIcon:
    def test_data_uri_adds_picture(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}image")
        el.set("href", _png_data_uri())
        el.set("transform", "matrix(1,0,0,1,50,60)")
        _process_icon(el, slide)
        assert len(slide.shapes) == 1

    def test_external_url_skipped(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}image")
        el.set("href", "https://example.com/img.png")
        el.set("transform", "matrix(1,0,0,1,0,0)")
        _process_icon(el, slide)
        assert len(slide.shapes) == 0

    def test_xlink_href_resolved(self):
        slide = _blank_slide()
        el = ET.Element(f"{{{SVG_NS}}}image")
        el.set(f"{{{XLINK_NS}}}href", _png_data_uri())
        el.set("transform", "matrix(1,0,0,1,0,0)")
        _process_icon(el, slide)
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# TestWalk
# ---------------------------------------------------------------------------

class TestWalk:
    def _svg(self):
        el = ET.Element(f"{{{SVG_NS}}}svg")
        el.set("width", "400px")
        el.set("height", "300px")
        return el

    def test_defs_skipped(self):
        svg = self._svg()
        ET.SubElement(svg, f"{{{SVG_NS}}}defs")
        slide = _blank_slide()
        _walk(svg, {}, {}, {}, slide, 400.0, 300.0)
        assert len(slide.shapes) == 0

    def test_top_level_rect_skipped(self):
        svg = self._svg()
        rect = ET.SubElement(svg, f"{{{SVG_NS}}}rect")
        rect.set("fill", "#ffffff")
        slide = _blank_slide()
        _walk(svg, {}, {}, {}, slide, 400.0, 300.0)
        assert len(slide.shapes) == 0

    def test_line_processed(self):
        svg = self._svg()
        line = ET.SubElement(svg, f"{{{SVG_NS}}}line")
        for attr, val in [("x1", "0"), ("y1", "0"), ("x2", "10"), ("y2", "10")]:
            line.set(attr, val)
        slide = _blank_slide()
        _walk(svg, {}, {}, {}, slide, 400.0, 300.0)
        assert len(slide.shapes) == 1

    def test_nested_group_processed(self):
        svg = self._svg()
        g = ET.SubElement(svg, f"{{{SVG_NS}}}g")
        line = ET.SubElement(g, f"{{{SVG_NS}}}line")
        for attr, val in [("x1", "0"), ("y1", "0"), ("x2", "5"), ("y2", "5")]:
            line.set(attr, val)
        slide = _blank_slide()
        _walk(svg, {}, {}, {}, slide, 400.0, 300.0)
        assert len(slide.shapes) == 1

    def test_multiple_elements_all_added(self):
        svg = self._svg()
        for i in range(3):
            line = ET.SubElement(svg, f"{{{SVG_NS}}}line")
            for attr, val in [("x1", "0"), ("y1", "0"),
                               ("x2", str(i + 1)), ("y2", "0")]:
                line.set(attr, val)
        slide = _blank_slide()
        _walk(svg, {}, {}, {}, slide, 400.0, 300.0)
        assert len(slide.shapes) == 3


# ---------------------------------------------------------------------------
# TestMain (integration — requires diagram_bricks.svg)
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not DIAGRAM_SVG.exists(), reason="diagram_bricks.svg not found")
class TestMain:
    @pytest.fixture(scope="class")
    def pptx_path(self, tmp_path_factory):
        out = tmp_path_factory.mktemp("pptx") / "out.pptx"
        with patch.object(sys, "argv", ["svg2pptx", str(DIAGRAM_SVG), "-o", str(out)]):
            main()
        return out

    def test_file_created(self, pptx_path):
        assert pptx_path.exists()

    def test_has_shapes(self, pptx_path):
        prs = Presentation(str(pptx_path))
        assert len(prs.slides[0].shapes) > 0

    def test_has_many_shapes(self, pptx_path):
        """diagram_bricks.svg has hundreds of brick <use> elements."""
        prs = Presentation(str(pptx_path))
        assert len(prs.slides[0].shapes) > 100

    def test_slide_dimensions_positive(self, pptx_path):
        prs = Presentation(str(pptx_path))
        assert prs.slide_width > 0
        assert prs.slide_height > 0
